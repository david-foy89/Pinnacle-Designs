"""Generate Pinnacle Designs investor/client pitch deck (.pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pitch-deck" / "pinnacle-designs-pitch-deck.pptx"
LOGO = ROOT / "assets" / "logo.png"

# Brand palette (from styles.css)
BLACK = RGBColor(0, 0, 0)
BG_ELEVATED = RGBColor(10, 10, 15)
BG_CARD = RGBColor(15, 16, 24)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(161, 161, 170)
ACCENT = RGBColor(59, 130, 246)
ACCENT_BRIGHT = RGBColor(96, 165, 250)
ACCENT_DEEP = RGBColor(37, 99, 235)
PERIWINKLE = RGBColor(129, 140, 248)


def set_slide_bg(slide, color: RGBColor = BLACK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top: float = 0, height: float = 0.06) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(top),
        Inches(13.333),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_DEEP
    bar.line.fill.background()


def add_footer(slide, text: str = "pinnacle-designs.com") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_label(slide, text: str, left=0.75, top=0.55) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BRIGHT
    p.font.name = "Calibri"


def add_title(slide, text: str, left=0.75, top=0.95, width=11.5, size=36) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_body(slide, text: str, left=0.75, top=2.0, width=11.5, size=18, color=MUTED) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.line_spacing = 1.3


def add_bullets(
    slide,
    items: list[str],
    left=0.75,
    top=2.1,
    width=11.5,
    size=17,
    bold_leads: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bold_leads and ": " in item:
            lead, rest = item.split(": ", 1)
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.level = 0
            p.space_after = Pt(10)
            # Rebuild with bold lead via runs
            p.text = ""
            r1 = p.add_run()
            r1.text = lead + ": "
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = WHITE
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = MUTED
            r2.font.name = "Calibri"
        else:
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = MUTED
            p.font.name = "Calibri"
            p.level = 0
            p.space_after = Pt(10)


def add_two_column_compare(
    slide,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    # Left column (muted / old way)
    card_l = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(2.0),
        Inches(5.75),
        Inches(4.6),
    )
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = BG_CARD
    card_l.line.color.rgb = RGBColor(80, 80, 90)

    tf_l = card_l.text_frame
    tf_l.margin_left = Inches(0.25)
    tf_l.margin_top = Inches(0.2)
    p = tf_l.paragraphs[0]
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED
    for item in left_items:
        bp = tf_l.add_paragraph()
        bp.text = f"✗  {item}"
        bp.font.size = Pt(14)
        bp.font.color.rgb = MUTED
        bp.space_before = Pt(8)

    # Right column (accent / pinnacle way)
    card_r = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.85),
        Inches(2.0),
        Inches(5.75),
        Inches(4.6),
    )
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = BG_ELEVATED
    card_r.line.color.rgb = ACCENT

    tf_r = card_r.text_frame
    tf_r.margin_left = Inches(0.25)
    tf_r.margin_top = Inches(0.2)
    p = tf_r.paragraphs[0]
    p.text = right_title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_BRIGHT
    for item in right_items:
        bp = tf_r.add_paragraph()
        bp.text = f"✓  {item}"
        bp.font.size = Pt(14)
        bp.font.color.rgb = WHITE
        bp.space_before = Pt(8)


def add_pricing_card(
    slide,
    left: float,
    tier: str,
    audience: str,
    down: str,
    monthly: str,
    features: list[str],
    featured: bool = False,
) -> None:
    w = 3.85
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(2.0),
        Inches(w),
        Inches(4.7),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_ELEVATED if featured else BG_CARD
    card.line.color.rgb = ACCENT if featured else RGBColor(60, 60, 70)
    card.line.width = Pt(2 if featured else 1)

    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = [tier, audience, f"{down} down", f"+ {monthly}/mo", ""] + [f"• {f}" for f in features]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = ACCENT_BRIGHT if featured else WHITE
        elif i == 1:
            p.font.size = Pt(11)
            p.font.color.rgb = MUTED
        elif i in (2, 3):
            p.font.bold = True
            p.font.size = Pt(16 if i == 2 else 14)
            p.font.color.rgb = WHITE
        else:
            p.font.size = Pt(11)
            p.font.color.rgb = MUTED
            p.space_before = Pt(2)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=6.95, height=0.45)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.15), Inches(1.2), width=Inches(5.0))

    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5))
    p = tagline.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "WEB DESIGN. "
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "ELEVATED."
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = ACCENT_BRIGHT
    r2.font.name = "Calibri"

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "Website-as-a-Service & industry management software for East Tennessee small businesses"
    sp.font.size = Pt(16)
    sp.font.color.rgb = MUTED
    sp.font.name = "Calibri"

    loc = slide.shapes.add_textbox(Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.4))
    lp = loc.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lp.text = "Erwin, TN · Tri-Cities · 2026"
    lp.font.size = Pt(12)
    lp.font.color.rgb = MUTED


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "The Problem")
    add_title(slide, "Local businesses are losing customers online")
    add_bullets(
        slide,
        [
            "Many shops rely only on Facebook — no owned digital storefront",
            "Outdated websites hurt credibility and fail on mobile",
            "Traditional agencies charge $3,000–$10,000+ upfront",
            "Owners are left to update hosting, security, and content themselves",
            "Generic software stacks don't fit how local industries actually operate",
        ],
        top=2.15,
        size=18,
    )
    add_footer(slide)


def slide_opportunity(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Market Opportunity")
    add_title(slide, "East Tennessee's underserved small-business market", size=32)
    add_bullets(
        slide,
        [
            "Target: plumbers, HVAC, roofers, restaurants, boutiques, and professional services",
            "Service area: Erwin, Unicoi County, Johnson City, Jonesborough, Kingsport",
            "High-intent local search: “web design Erwin TN,” “affordable web developer Johnson City”",
            "WaaS model lowers barrier to entry and creates predictable recurring revenue",
            "Expansion path: industry-specific management software beyond websites",
        ],
        top=2.15,
        bold_leads=True,
    )
    add_footer(slide)


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Our Solution")
    add_title(slide, "Website-as-a-Service (WaaS)", size=34)
    add_body(
        slide,
        "Pinnacle Designs builds conversion-focused websites for local businesses — then acts as the outsourced webmaster: hosting, SSL, security, SEO, and done-for-you updates every month.",
        top=1.95,
        size=19,
        color=WHITE,
    )
    add_bullets(
        slide,
        [
            "Small down payment + flat monthly subscription",
            "Text us when prices or photos change — we handle the rest",
            "Continuous maintenance so sites never go stale in ~2 years",
            "Zero tech headaches for the business owner",
        ],
        top=3.35,
        size=17,
    )
    add_footer(slide)


def slide_compare(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Why Pinnacle")
    add_title(slide, "The old way vs. the Pinnacle way", size=32)
    add_two_column_compare(
        slide,
        "Traditional Agencies",
        [
            "Huge $3,000+ upfront cost",
            "You update the site yourself",
            "Pay extra for hosting & security",
            "Site goes out of date in ~2 years",
        ],
        "Pinnacle Designs",
        [
            "Affordable down payment from $500",
            "We make updates for you",
            "Hosting & security included",
            "Continuously maintained & modern",
        ],
    )
    add_footer(slide)


def slide_pricing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Services & Pricing")
    add_title(slide, "Three plans. One simple model.", size=32)
    add_pricing_card(
        slide,
        0.75,
        "Base Camp",
        "New businesses & contractors",
        "$500",
        "$99",
        ["3–5 page site", "Basic SEO", "Hosting & SSL", "1 hr/mo updates"],
    )
    add_pricing_card(
        slide,
        4.75,
        "Ascent",
        "Established local shops",
        "$1,000",
        "$199",
        ["Up to 10 pages", "Google Business Profile", "Analytics", "2 hrs/mo updates"],
        featured=True,
    )
    add_pricing_card(
        slide,
        8.75,
        "Summit",
        "Retailers selling online",
        "$2,500+",
        "$299+",
        ["Full store setup", "Payment gateways", "Priority support", "Catalog updates"],
    )
    note = slide.shapes.add_textbox(Inches(0.75), Inches(6.85), Inches(12), Inches(0.35))
    np = note.text_frame.paragraphs[0]
    np.text = "Design fee + monthly retainer. Larger add-ons quoted separately in contract."
    np.font.size = Pt(11)
    np.font.color.rgb = MUTED
    add_footer(slide)


def slide_value(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Client Outcomes")
    add_title(slide, "What clients get", size=34)

    outcomes = [
        ("Never worry about hacks", "Security scanning and managed hosting — protected without thinking about it."),
        ("Show up locally", "Mobile-friendly design and SEO for Unicoi County and the Tri-Cities."),
        ("Unlimited peace of mind", "Monthly updates keep menus, services, and promotions current."),
    ]
    for i, (head, body) in enumerate(outcomes):
        left = 0.75 + i * 4.1
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.1),
            Inches(3.85),
            Inches(3.2),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = ACCENT

        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = "◆"
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_BRIGHT
        h = tf.add_paragraph()
        h.text = head
        h.font.bold = True
        h.font.size = Pt(16)
        h.font.color.rgb = WHITE
        b = tf.add_paragraph()
        b.text = body
        b.font.size = Pt(12)
        b.font.color.rgb = MUTED

    guarantee = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(5.55),
        Inches(11.85),
        Inches(0.9),
    )
    guarantee.fill.solid()
    guarantee.fill.fore_color.rgb = BG_ELEVATED
    guarantee.line.color.rgb = PERIWINKLE
    gp = guarantee.text_frame.paragraphs[0]
    gp.text = "If you don't love the design we present, you don't pay the down payment. Zero risk to get started."
    gp.font.size = Pt(15)
    gp.font.color.rgb = WHITE
    gp.alignment = PP_ALIGN.CENTER
    guarantee.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(slide)


def slide_software(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Growth — Management Software")
    add_title(slide, "Beyond websites: tools built per industry", size=30)
    add_body(
        slide,
        "In active development — industry-specific management software so local businesses handle jobs, customers, and operations without juggling a dozen apps.",
        top=1.85,
        size=17,
        color=MUTED,
    )

    verticals = [
        ("Home Services", "Scheduling, customer records, follow-ups"),
        ("Retail & Boutiques", "Inventory, customers, sales tracking"),
        ("Restaurants", "Menus, order flow, staff coordination"),
        ("Professional Services", "Intake, booking, project tracking"),
    ]
    for i, (name, desc) in enumerate(verticals):
        row, col = divmod(i, 2)
        left = 0.75 + col * 6.1
        top = 2.85 + row * 1.85
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(5.75),
            Inches(1.55),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = RGBColor(50, 50, 60)
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = name
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_BRIGHT
        d = tf.add_paragraph()
        d.text = desc
        d.font.size = Pt(12)
        d.font.color.rgb = MUTED

    badge = slide.shapes.add_textbox(Inches(0.75), Inches(6.55), Inches(5), Inches(0.35))
    bp = badge.text_frame.paragraphs[0]
    bp.text = "IN DEVELOPMENT — Early access waitlist open"
    bp.font.bold = True
    bp.font.size = Pt(11)
    bp.font.color.rgb = ACCENT_BRIGHT
    add_footer(slide)


def slide_business_model(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Business Model")
    add_title(slide, "Predictable revenue, low client friction", size=32)
    add_bullets(
        slide,
        [
            "Revenue stream 1: One-time design & setup fees ($500 – $2,500+)",
            "Revenue stream 2: Recurring monthly management ($99 – $299+/mo)",
            "Revenue stream 3: Out-of-scope projects (new pages, integrations, e-commerce add-ons)",
            "Revenue stream 4 (future): Management software subscriptions by vertical",
            "Low CAC: Local SEO, referrals, Facebook, and in-person Tri-Cities networking",
            "High retention: Ongoing updates create sticky, long-term client relationships",
        ],
        top=2.1,
        bold_leads=True,
        size=16,
    )
    add_footer(slide)


def slide_traction(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "Traction & Readiness")
    add_title(slide, "Built to scale locally", size=34)

    stats = [
        ("$500", "Starting down payment"),
        ("24/7", "Digital storefront for clients"),
        ("0", "Tech headaches for owners"),
    ]
    for i, (val, label) in enumerate(stats):
        left = 0.75 + i * 4.1
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.2),
            Inches(3.85),
            Inches(2.0),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_ELEVATED
        card.line.color.rgb = ACCENT
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = val
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = ACCENT_BRIGHT
        l = tf.add_paragraph()
        l.alignment = PP_ALIGN.CENTER
        l.text = label
        l.font.size = Pt(13)
        l.font.color.rgb = MUTED

    add_bullets(
        slide,
        [
            "Live brand & marketing site at pinnacle-designs.com",
            "Contract templates ready for Base Camp, Ascent, Summit, and software beta",
            "Formspree contact pipeline and Google Analytics in place",
            "Cloudflare-hosted with SEO, FAQ schema, and local landing content",
        ],
        top=4.55,
        size=16,
    )
    add_footer(slide)


def slide_team_ask(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_label(slide, "The Ask")
    add_title(slide, "Partner with Pinnacle Designs", size=34)
    add_bullets(
        slide,
        [
            "For clients: Start with a free conversation — pick Base Camp, Ascent, or Summit",
            "For software beta: Request early access for your industry vertical",
            "For partners & investors: Let's talk about scaling WaaS and software across East TN",
        ],
        top=2.1,
        size=18,
    )

    contact = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(4.35),
        Inches(11.85),
        Inches(2.0),
    )
    contact.fill.solid()
    contact.fill.fore_color.rgb = BG_ELEVATED
    contact.line.color.rgb = ACCENT
    tf = contact.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = [
        "david@pinnacle-designs.com",
        "pinnacle-designs.com",
        "facebook.com/profile.php?id=61590691171007",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        p.font.size = Pt(20 if i == 0 else 16)
        p.font.color.rgb = ACCENT_BRIGHT if i == 0 else MUTED
        p.font.bold = i == 0

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(5.4), Inches(6.55), width=Inches(2.5))

    add_footer(slide, "")


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_opportunity(prs)
    slide_solution(prs)
    slide_compare(prs)
    slide_pricing(prs)
    slide_value(prs)
    slide_software(prs)
    slide_business_model(prs)
    slide_traction(prs)
    slide_team_ask(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Created {path}")
